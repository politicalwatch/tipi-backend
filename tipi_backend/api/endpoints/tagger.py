import codecs
import logging
import pickle
import subprocess
from pdfminer.high_level import extract_text as extract_pdf_text
from docx import Document
from pptx import Presentation
from os.path import splitext
import tempfile

from flask import request, abort
from flask_restx import Namespace, Resource

import tipi_tasks
from tipi_backend.api.business import get_tags, get_kbs
from tipi_backend.api.endpoints import cache, limiter
from tipi_backend.api.parsers import parser_tagger, parser_kb
from tipi_backend.settings import Config


log = logging.getLogger(__name__)

ns = Namespace(
    "tagger", description="Operations related to tag texts using our knowledge base"
)


def filter_tags(result, kb):
    topics = result["result"]["topics"]
    tags = result["result"]["tags"]
    new_topics = []
    new_tags = []
    for tag in tags:
        if tag["knowledgebase"] in kb:
            new_tags.append(tag)
            new_topics.append(tag["topic"])
    new_topics = list(set(new_topics))
    result["result"]["topics"] = new_topics
    result["result"]["tags"] = new_tags

    return result


def remove_fields(result):
    tags = result["result"]["tags"]
    for tag in tags:
        del tag["public"]


@ns.route("/")
@ns.expect(parser_tagger)
class TaggerExtractor(Resource):

    def post(self):
        """Returns a list of topics and tags matching the text."""
        try:
            args = parser_tagger.parse_args(request)
            kb = get_kbs(args)

            cache_key = Config.CACHE_TAGS
            tags = cache.get(cache_key)
            if tags is None:
                tags = get_tags()
                cache.set(cache_key, tags, timeout=5 * 60)
            tags = codecs.encode(pickle.dumps(tags), "base64").decode()
            tipi_tasks.init()
            text = ""
            if "text" in args and args["text"]:
                text = args["text"]
            else:
                if "file" in args:
                    file_input = args["file"]
                    with tempfile.NamedTemporaryFile(
                        prefix="tipiscanner_", suffix=splitext(file_input.filename)[1]
                    ) as f:
                        f.write(file_input.stream.read())
                        f.seek(0)
                        print("MIMETYPE:", file_input.mimetype)
                        if file_input.mimetype == "text/plain":
                            text = f.read().decode("utf-8").strip()
                        elif file_input.mimetype == "application/pdf":
                            text = extract_pdf_text(f.name).strip()
                        elif (
                            file_input.mimetype
                            == "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
                        ):
                            doc = Document(f)
                            text = "\n".join(
                                [para.text for para in doc.paragraphs]
                            ).strip()
                        elif file_input.mimetype == "application/msword":
                            result = subprocess.run(
                                ["antiword", f.name],
                                stdout=subprocess.PIPE,
                                stderr=subprocess.PIPE,
                            )
                            if result.returncode != 0:
                                raise Exception(
                                    f"Error al leer el archivo .doc: {result.stderr.decode('utf-8')}"
                                )
                            text = result.stdout.decode("utf-8").strip()
                        elif (
                            file_input.mimetype
                            == "application/vnd.openxmlformats-officedocument.presentationml.presentation"
                        ):
                            ppt = Presentation(f)
                            text = "\n".join(
                                [
                                    shape.text
                                    for slide in ppt.slides
                                    for shape in slide.shapes
                                    if hasattr(shape, "text")
                                ]
                            ).strip()
                        else:
                            abort(
                                400,
                                "Formato no soportado. Por favor, utilice un archivo .txt, .pdf, .docx, .doc o .pptx.",
                            )
                        f.close()
                    if not text:
                        abort(
                            400,
                            "Error al obtener el texto del fichero proporcionado. Pruebe con otro fichero.",
                        )
                    print("Texto:", text)
            text_length = len(text.split())

            if text_length >= Config.TAGGER_MAX_WORDS:
                task = tipi_tasks.tagger.extract_tags_from_text.apply_async(
                    (text, tags)
                )
                eta_time = int((text_length / 1000) * 4)
                task_id = task.id
                result = {
                    "status": "PROCESSING",
                    "task_id": task_id,
                    "estimated_time": eta_time,
                }
            else:
                result = tipi_tasks.tagger.extract_tags_from_text(text, tags)
                result = filter_tags(result, kb)
                remove_fields(result)

            return result
        except Exception as e:
            if hasattr(e, "code") and hasattr(e, "description"):
                abort(e.code, e.description)
            else:
                abort(500, "Internal server error")


@ns.route("/result/<id>")
@ns.param(
    name="id",
    description="Task id",
    type=str,
    required=True,
    location=["path"],
    help="Invalid identifier",
)
@ns.response(404, "Task not found.")
@ns.expect(parser_kb)
class TaggerResult(Resource):

    def get(self, id):
        """Returns tagging task's result"""
        try:
            tipi_tasks.init()
            result = tipi_tasks.tagger.check_status_task(id)

            args = parser_kb.parse_args(request)
            kb = get_kbs(args)

            if result["status"] == "SUCCESS":
                result = filter_tags(result, kb)
                remove_fields(result)

            return result
        except Exception as e:
            log.error(e)
            return {"Error": "No task found"}, 404
